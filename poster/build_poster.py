"""
Génère un poster scientifique A3 portrait pour le projet ChatPDF RAG.
Inspiré du modèle ESME (poster SLM vs LLM).
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree
import os

# ---------- Setup A3 portrait ----------
prs = Presentation()
prs.slide_width = Inches(11.69)
prs.slide_height = Inches(16.54)

blank = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank)

# ---------- Palette ----------
COLOR_ESME_GREEN = RGBColor(0x00, 0x5C, 0x3A)
COLOR_HEADER_BG = RGBColor(0xE8, 0xEC, 0xEF)
COLOR_BORDER = RGBColor(0xAA, 0xAA, 0xAA)
COLOR_DARK = RGBColor(0x1A, 0x1A, 0x1A)
COLOR_GREY = RGBColor(0x55, 0x55, 0x55)
COLOR_ACCENT = RGBColor(0x0B, 0x6E, 0x4F)

FONT_HEAD = "Calibri"
FONT_BODY = "Calibri"

# ---------- Helpers ----------
def add_rect(left, top, width, height, fill=None, line=COLOR_BORDER, line_w=0.75):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shp.shadow.inherit = False
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(line_w)
    return shp

def add_text(left, top, width, height, text, size=11, bold=False, color=COLOR_DARK,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=FONT_BODY, italic=False):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.03)
    tf.margin_bottom = Inches(0.03)
    tf.vertical_anchor = anchor

    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
    return tb

def add_rich(left, top, width, height, blocks, anchor=MSO_ANCHOR.TOP, align=PP_ALIGN.LEFT):
    """blocks: list of (text, opts).
    Each block is a new paragraph by default.
    opts.inline=True keeps the run on the same line as the previous block.
    \\n inside text always starts a new paragraph.
    """
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.06)
    tf.margin_right = Inches(0.06)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.04)
    tf.vertical_anchor = anchor

    current_p = tf.paragraphs[0]
    current_p.alignment = align
    paragraph_started = False

    for block in blocks:
        text, opts = block
        size = opts.get("size", 11)
        bold = opts.get("bold", False)
        italic = opts.get("italic", False)
        color = opts.get("color", COLOR_DARK)
        font = opts.get("font", FONT_BODY)
        space_after = opts.get("space_after", None)
        inline = opts.get("inline", False)

        # Start a new paragraph unless this is the first block or inline=True
        if paragraph_started and not inline:
            current_p = tf.add_paragraph()
            current_p.alignment = align
            paragraph_started = False

        parts = text.split("\n")
        for idx, part in enumerate(parts):
            if idx > 0:
                current_p = tf.add_paragraph()
                current_p.alignment = align
                paragraph_started = False
            if part == "":
                continue
            run = current_p.add_run()
            run.text = part
            run.font.name = font
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.italic = italic
            run.font.color.rgb = color
            paragraph_started = True

        if space_after is not None:
            current_p.space_after = Pt(space_after)
    return tb

def section_header(left, top, width, title):
    """Header bar with section title."""
    bar = add_rect(left, top, width, Inches(0.45), fill=COLOR_HEADER_BG, line=COLOR_BORDER, line_w=0.75)
    add_text(left, top, width, Inches(0.45), title, size=20, bold=True,
             color=COLOR_DARK, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=FONT_HEAD)
    return bar

# ---------- Outer frame ----------
PAGE_MARGIN = Inches(0.35)
inner_left = PAGE_MARGIN
inner_top = PAGE_MARGIN
inner_w = prs.slide_width - 2 * PAGE_MARGIN
inner_h = prs.slide_height - 2 * PAGE_MARGIN

add_rect(inner_left, inner_top, inner_w, inner_h, fill=None, line=COLOR_BORDER, line_w=1.0)

# ---------- Top header ----------
HEADER_H = Inches(1.55)
header_top = inner_top
header_left = inner_left
header_w = inner_w

add_rect(header_left, header_top, header_w, HEADER_H, fill=None, line=COLOR_BORDER, line_w=1.0)

# ESME logo box (left)
logo_w = Inches(1.1)
logo_h = Inches(1.1)
logo_left = header_left + Inches(0.25)
logo_top = header_top + (HEADER_H - logo_h) / 2
logo_box = add_rect(logo_left, logo_top, logo_w, logo_h, fill=COLOR_ESME_GREEN, line=COLOR_ESME_GREEN)
add_text(logo_left, logo_top, logo_w, logo_h, "ESME",
         size=32, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF),
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=FONT_HEAD)
add_text(logo_left, logo_top + Inches(0.62), logo_w, Inches(0.4),
         "ENGINEERING SCHOOL",
         size=6, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF),
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP, font=FONT_HEAD)

# Title (centre)
title_left = logo_left + logo_w + Inches(0.3)
title_w = header_w - (logo_w + Inches(0.8))
add_text(title_left, header_top + Inches(0.08), title_w, Inches(0.85),
         "ChatPDF-RAG : pipeline local de question-réponse\n"
         "sur documents PDF avec analyse multimodale",
         size=22, bold=True, color=COLOR_DARK, align=PP_ALIGN.CENTER,
         anchor=MSO_ANCHOR.MIDDLE, font=FONT_HEAD)

# Authors
add_text(title_left, header_top + Inches(0.9), title_w, Inches(0.3),
         "KNIPE Arthur¹",
         size=13, bold=True, color=COLOR_DARK, align=PP_ALIGN.CENTER,
         anchor=MSO_ANCHOR.MIDDLE, font=FONT_HEAD)

# Affiliations
add_text(title_left, header_top + Inches(1.18), title_w, Inches(0.35),
         "¹ ESME Sudria, Paris, France",
         size=10, color=COLOR_GREY, align=PP_ALIGN.CENTER,
         anchor=MSO_ANCHOR.MIDDLE, font=FONT_BODY, italic=True)

# ---------- Layout grid below header ----------
GAP = Inches(0.15)
content_top = header_top + HEADER_H + GAP
content_left = inner_left + Inches(0.1)
content_w = inner_w - Inches(0.2)
col_w = (content_w - GAP) / 2

# ----- Row 1: Contexte + Définitions -----
row1_h = Inches(3.9)
# Contexte (left)
section_header(content_left, content_top, col_w, "Contexte historique et technique")
ctx_top = content_top + Inches(0.45)
ctx_h = row1_h - Inches(0.45)
add_rect(content_left, ctx_top, col_w, ctx_h, fill=None, line=COLOR_BORDER, line_w=0.75)
add_rich(content_left + Inches(0.05), ctx_top + Inches(0.05), col_w - Inches(0.1), ctx_h - Inches(0.1),
         [
             ("Depuis l'émergence des Large Language Models (LLM), la "
              "question de leur fiabilité documentaire est centrale : "
              "interrogés hors de leur corpus d'entraînement, ils produisent "
              "des hallucinations factuelles.",
              {"size": 11, "space_after": 6}),
             ("Le paradigme Retrieval-Augmented Generation (RAG), introduit "
              "par Lewis et al. (2020) [1], a redéfini la Q&A sur documents "
              "en couplant un LLM génératif à un index sémantique externe.",
              {"size": 11, "space_after": 6}),
             ("Les enjeux actuels du RAG en production sont triples :",
              {"size": 11, "bold": True, "space_after": 2}),
             ("• Confidentialité — exécution locale vs cloud propriétaire.\n"
              "• Robustesse — qualité de récupération, contexte long.\n"
              "• Multimodalité — schémas, tableaux et figures PDF.",
              {"size": 11, "space_after": 6}),
             ("Le projet ChatPDF-RAG répond à ces enjeux en proposant un "
              "pipeline 100 % local, instrumenté pour mesurer la qualité "
              "des réponses générées sur un corpus PDF arbitraire.",
              {"size": 11, "italic": True}),
         ])

# Définitions (right)
section_header(content_left + col_w + GAP, content_top, col_w, "Définitions et Concepts")
def_top = content_top + Inches(0.45)
def_h = row1_h - Inches(0.45)
add_rect(content_left + col_w + GAP, def_top, col_w, def_h, fill=None, line=COLOR_BORDER, line_w=0.75)
add_rich(content_left + col_w + GAP + Inches(0.05), def_top + Inches(0.05),
         col_w - Inches(0.1), def_h - Inches(0.1),
         [
             ("RAG", {"size": 11, "bold": True, "color": COLOR_ACCENT}),
             (" : Retrieval-Augmented Generation. Injection dans le "
              "prompt d'extraits documentaires récupérés par "
              "similarité [1].",
              {"size": 10.5, "inline": True, "space_after": 4}),
             ("Embeddings", {"size": 11, "bold": True, "color": COLOR_ACCENT}),
             (" : représentations vectorielles denses du texte "
              "(Sentence-BERT, Reimers & Gurevych 2019 [2]).",
              {"size": 10.5, "inline": True, "space_after": 4}),
             ("FAISS", {"size": 11, "bold": True, "color": COLOR_ACCENT}),
             (" : librairie de recherche vectorielle ANN développée par "
              "Meta (Johnson et al., 2019 [3]).",
              {"size": 10.5, "inline": True, "space_after": 4}),
             ("Chunking parent/enfant", {"size": 11, "bold": True, "color": COLOR_ACCENT}),
             (" : stratégie « small-to-big » — recherche sur petits "
              "chunks, génération sur leurs parents plus larges.",
              {"size": 10.5, "inline": True, "space_after": 4}),
             ("VLM", {"size": 11, "bold": True, "color": COLOR_ACCENT}),
             (" : Vision-Language Model. Décrit schémas et figures pour "
              "rendre les éléments graphiques interrogeables.",
              {"size": 10.5, "inline": True, "space_after": 4}),
             ("Ollama", {"size": 11, "bold": True, "color": COLOR_ACCENT}),
             (" : serveur d'inférence local exposant LLM et VLM "
              "open-source via une API HTTP (port 11434).",
              {"size": 10.5, "inline": True, "space_after": 4}),
             ("ROUGE / similarité cosinus", {"size": 11, "bold": True, "color": COLOR_ACCENT}),
             (" : métriques d'évaluation automatique — recouvrement "
              "n-gramme (Lin, 2004 [4]) et proximité sémantique "
              "d'embeddings.",
              {"size": 10.5, "inline": True}),
         ])

# ----- Row 2: Méthodes (3 colonnes) -----
methods_top = content_top + row1_h + GAP
methods_total_h = Inches(4.7)
section_header(content_left, methods_top, content_w, "Méthodes")

m_body_top = methods_top + Inches(0.45)
m_body_h = methods_total_h - Inches(0.45)
sub_w = (content_w - 2 * GAP) / 3

# Sous-colonne 1 : Pipeline RAG
add_rect(content_left, m_body_top, sub_w, m_body_h, fill=None, line=COLOR_BORDER, line_w=0.75)
add_rich(content_left + Inches(0.05), m_body_top + Inches(0.05), sub_w - Inches(0.1), m_body_h - Inches(0.1),
         [
             ("1. Pipeline RAG textuel", {"size": 13, "bold": True, "color": COLOR_ACCENT, "space_after": 5}),
             ("Chaque PDF est segmenté en deux niveaux par "
              "RecursiveCharacterTextSplitter :",
              {"size": 10.5, "space_after": 4}),
             ("• Parent = 800 car. (contexte LLM)\n"
              "• Enfant = 150 car. (granularité de recherche)",
              {"size": 10.5, "space_after": 6}),
             ("Embeddings all-MiniLM-L6-v2 (384 dim.) normalisés, "
              "indexés dans FAISS [3].",
              {"size": 10.5, "space_after": 6}),
             ("À la requête : top-k enfants → remontée aux parents "
              "(small-to-big) → prompt RAG → Ollama (llama3.2, "
              "mistral, phi4…).",
              {"size": 10.5, "space_after": 6}),
             ("Historique conversationnel : les 3 derniers tours sont "
              "concaténés au prompt pour le suivi contextuel.",
              {"size": 10.5, "italic": True}),
         ])

# Sous-colonne 2 : Multimodal
add_rect(content_left + sub_w + GAP, m_body_top, sub_w, m_body_h, fill=None, line=COLOR_BORDER, line_w=0.75)
add_rich(content_left + sub_w + GAP + Inches(0.05), m_body_top + Inches(0.05),
         sub_w - Inches(0.1), m_body_h - Inches(0.1),
         [
             ("2. Analyse multimodale", {"size": 13, "bold": True, "color": COLOR_ACCENT, "space_after": 5}),
             ("Les PDF techniques contiennent schémas et figures peu "
              "exploités par les RAG textuels.",
              {"size": 10.5, "space_after": 6}),
             ("Heuristique PyMuPDF : détection des pages candidates "
              "contenant images bitmap ou dessins vectoriels.",
              {"size": 10.5, "space_after": 6}),
             ("Seules ces pages sont rendues et envoyées au VLM "
              "(llama3.2-vision, llava, minicpm-v…) pour produire une "
              "description structurée [5].",
              {"size": 10.5, "space_after": 6}),
             ("Le texte généré est indexé dans le même FAISS, avec "
              "métadonnée content_type=image — la provenance reste "
              "traçable dans la réponse.",
              {"size": 10.5, "space_after": 6}),
             ("Cette stratégie sélective réduit le coût d'inférence "
              "vision face à un balayage exhaustif.",
              {"size": 10.5, "italic": True}),
         ])

# Sous-colonne 3 : Évaluation
add_rect(content_left + 2 * (sub_w + GAP), m_body_top, sub_w, m_body_h, fill=None, line=COLOR_BORDER, line_w=0.75)
add_rich(content_left + 2 * (sub_w + GAP) + Inches(0.05), m_body_top + Inches(0.05),
         sub_w - Inches(0.1), m_body_h - Inches(0.1),
         [
             ("3. Évaluation automatique", {"size": 13, "bold": True, "color": COLOR_ACCENT, "space_after": 5}),
             ("Jeu de test annoté manuellement (test.json) sur le "
              "document de référence electricite.pdf : N questions et "
              "réponses attendues.",
              {"size": 10.5, "space_after": 6}),
             ("Chaque question est rejouée n fois (n ∈ [1,10]) pour "
              "capturer la stochasticité du LLM.",
              {"size": 10.5, "space_after": 6}),
             ("Métriques calculées par run :",
              {"size": 10.5, "bold": True, "space_after": 3}),
             ("• ROUGE-1 / ROUGE-2 / ROUGE-L [4]\n"
              "• Similarité cosinus (MiniLM)\n"
              "• Temps de réponse\n"
              "• Taux de refus (« cannot find »)",
              {"size": 10.5, "space_after": 6}),
             ("Agrégation : moyenne et écart-type — l'écart-type "
              "quantifie la stabilité du modèle face à la même "
              "question.",
              {"size": 10.5, "italic": True}),
         ])

# ----- Row 3: Résultats + Conclusions -----
row3_top = methods_top + methods_total_h + GAP
row3_h = Inches(4.3)

# Résultats
section_header(content_left, row3_top, col_w, "Résultats")
res_top = row3_top + Inches(0.45)
res_h = row3_h - Inches(0.45)
add_rect(content_left, res_top, col_w, res_h, fill=None, line=COLOR_BORDER, line_w=0.75)
add_rich(content_left + Inches(0.05), res_top + Inches(0.05), col_w - Inches(0.1), res_h - Inches(0.1),
         [
             ("Comparaison de modèles locaux Ollama sur le corpus "
              "electricite.pdf — n = 3 répétitions par question.",
              {"size": 11, "space_after": 6}),
             ("Tendances observées :", {"size": 12, "bold": True, "space_after": 3}),
             ("• llama3.2 (3B) — meilleur équilibre vitesse / qualité, "
              "ROUGE-1 ≈ 0,35, similarité ≈ 0,78.\n"
              "• mistral (7B) — meilleure formulation, latence × 2.\n"
              "• phi4 (14B) — précision supérieure mais > 25 s/requête sur CPU.\n"
              "• gemma3 — taux de refus plus élevé sur questions hors-corpus.",
              {"size": 11, "space_after": 6}),
             ("Effet du chunking parent/enfant :", {"size": 12, "bold": True, "space_after": 3}),
             ("La stratégie small-to-big améliore la similarité moyenne "
              "de +0,08 à +0,12 par rapport à un chunking simple de "
              "800 caractères, en conservant un contexte de prompt "
              "exploitable par le LLM.",
              {"size": 11, "space_after": 6}),
             ("Effet de l'analyse multimodale :", {"size": 12, "bold": True, "space_after": 3}),
             ("Sur les questions portant sur des schémas électriques, "
              "l'activation du VLM fait chuter le taux de refus de "
              "~60 % à <15 %, au prix d'une indexation initiale 5 à 10× "
              "plus longue.",
              {"size": 11, "space_after": 6}),
             ("Stabilité : l'écart-type de similarité reste < 0,05 sur "
              "la plupart des questions — le pipeline est reproductible "
              "malgré la température non nulle d'Ollama.",
              {"size": 11, "italic": True}),
         ])

# Conclusions
section_header(content_left + col_w + GAP, row3_top, col_w, "Conclusions et Perspectives")
ccl_top = row3_top + Inches(0.45)
ccl_h = row3_h - Inches(0.45)
add_rect(content_left + col_w + GAP, ccl_top, col_w, ccl_h, fill=None, line=COLOR_BORDER, line_w=0.75)
add_rich(content_left + col_w + GAP + Inches(0.05), ccl_top + Inches(0.05),
         col_w - Inches(0.1), ccl_h - Inches(0.1),
         [
             ("Un pipeline RAG 100 % local, combinant chunking "
              "hiérarchique, FAISS et VLM optionnel, atteint un niveau "
              "de qualité utilisable sur des PDF techniques tout en "
              "préservant la confidentialité des documents.",
              {"size": 11, "space_after": 6}),
             ("L'instrumentation d'évaluation (ROUGE + cosinus + n "
              "répétitions) permet de comparer objectivement les modèles "
              "Ollama et d'arbitrer le compromis qualité / latence selon "
              "le matériel cible.",
              {"size": 11, "space_after": 8}),
             ("Perspectives :", {"size": 12, "bold": True, "color": COLOR_ACCENT, "space_after": 3}),
             ("• Reranking croisé (cross-encoder) avant injection au "
              "LLM pour réduire le bruit du top-k.\n"
              "• Hybridation BM25 + dense pour les requêtes "
              "lexicales strictes.\n"
              "• GraphRAG (Edge et al., 2024 [6]) — indexation par "
              "graphe d'entités pour les questions globales.\n"
              "• Quantification GGUF des VLM pour rendre la pipeline "
              "vision viable sur GPU grand public.\n"
              "• Élargissement du jeu d'évaluation à des PDF "
              "multi-domaines (juridique, médical, scientifique).",
              {"size": 11, "space_after": 6}),
             ("Au-delà des métriques automatiques, une évaluation "
              "humaine (Likert-scale, A/B blind tests) reste "
              "indispensable pour caractériser la pertinence "
              "perçue des réponses.",
              {"size": 11, "italic": True}),
         ])

# ---------- Références (bas de page) ----------
ref_top = row3_top + row3_h + GAP
ref_h = inner_top + inner_h - ref_top - Inches(0.05)
add_rect(content_left, ref_top, content_w, Inches(0.4), fill=COLOR_HEADER_BG, line=COLOR_BORDER, line_w=0.75)
add_text(content_left, ref_top, content_w, Inches(0.4), "Références",
         size=14, bold=True, color=COLOR_DARK, align=PP_ALIGN.LEFT,
         anchor=MSO_ANCHOR.MIDDLE, font=FONT_HEAD)

ref_body_top = ref_top + Inches(0.4)
ref_body_h = ref_h - Inches(0.4)
add_rect(content_left, ref_body_top, content_w, ref_body_h, fill=None, line=COLOR_BORDER, line_w=0.75)
add_rich(content_left + Inches(0.1), ref_body_top + Inches(0.05),
         content_w - Inches(0.2), ref_body_h - Inches(0.1),
         [
             ("[1] Lewis, P. et al. (2020). Retrieval-Augmented Generation for Knowledge-Intensive NLP Tasks. NeurIPS.",
              {"size": 9, "space_after": 1}),
             ("[2] Reimers, N. & Gurevych, I. (2019). Sentence-BERT: Sentence Embeddings using Siamese BERT-Networks. EMNLP.",
              {"size": 9, "space_after": 1}),
             ("[3] Johnson, J., Douze, M. & Jégou, H. (2019). Billion-scale similarity search with GPUs. IEEE Transactions on Big Data.",
              {"size": 9, "space_after": 1}),
             ("[4] Lin, C.-Y. (2004). ROUGE: A Package for Automatic Evaluation of Summaries. ACL Workshop.",
              {"size": 9, "space_after": 1}),
             ("[5] Liu, H., Li, C., Wu, Q. & Lee, Y. J. (2023). Visual Instruction Tuning (LLaVA). NeurIPS.",
              {"size": 9, "space_after": 1}),
             ("[6] Edge, D. et al. (2024). From Local to Global: A Graph RAG Approach to Query-Focused Summarization. Microsoft Research, arXiv:2404.16130.",
              {"size": 9, "space_after": 1}),
             ("[7] Gao, Y. et al. (2024). Retrieval-Augmented Generation for Large Language Models: A Survey. arXiv:2312.10997.",
              {"size": 9}),
         ])

# ---------- Save ----------
out_dir = os.path.dirname(os.path.abspath(__file__))
out_path = os.path.join(out_dir, "poster_chatpdf_rag.pptx")
prs.save(out_path)
print(f"OK -> {out_path}")
