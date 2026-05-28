"""
Génère la slide "Métriques d'évaluation" qui manque dans
AI Music Industry.pdf (slide 6 du deck), au style de la slide 5
"Outils Utilisés" : titre serif italique, 4 colonnes bannière,
fond crème, blob bleu en haut à droite.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ---------- 16:9 widescreen ----------
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
slide = prs.slides.add_slide(prs.slide_layouts[6])

# ---------- Palette (matchée sur la slide de référence) ----------
COLOR_BG = RGBColor(0xF6, 0xF2, 0xE8)        # crème
COLOR_BLOB_LIGHT = RGBColor(0xD4, 0xDD, 0xEC) # bleu très clair
COLOR_BLOB_MID = RGBColor(0x9D, 0xB0, 0xD3)   # bleu mid
COLOR_DARK = RGBColor(0x1A, 0x1A, 0x1A)
COLOR_ACCENT = RGBColor(0x2E, 0x3F, 0x6E)
COLOR_LINE = RGBColor(0x1A, 0x1A, 0x1A)
COLOR_GREY = RGBColor(0x66, 0x66, 0x66)

FONT_TITLE = "Cambria"
FONT_BODY = "Calibri"

# ---------- Fond crème pleine page ----------
bg = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
)
bg.shadow.inherit = False
bg.fill.solid()
bg.fill.fore_color.rgb = COLOR_BG
bg.line.fill.background()

# ---------- Blob bleu décoratif en haut à droite ----------
blob_outer = slide.shapes.add_shape(
    MSO_SHAPE.OVAL,
    Inches(10.6), Inches(-1.2), Inches(4.5), Inches(3.2)
)
blob_outer.shadow.inherit = False
blob_outer.fill.solid()
blob_outer.fill.fore_color.rgb = COLOR_BLOB_LIGHT
blob_outer.line.fill.background()

blob_inner = slide.shapes.add_shape(
    MSO_SHAPE.OVAL,
    Inches(11.2), Inches(-0.6), Inches(3.0), Inches(2.0)
)
blob_inner.shadow.inherit = False
blob_inner.fill.solid()
blob_inner.fill.fore_color.rgb = COLOR_BLOB_MID
blob_inner.line.fill.background()

# Petite croix décorative (comme sur l'original)
cross_tb = slide.shapes.add_textbox(
    Inches(12.5), Inches(0.15), Inches(0.5), Inches(0.6)
)
ctf = cross_tb.text_frame
ctf.margin_left = ctf.margin_right = ctf.margin_top = ctf.margin_bottom = 0
p = ctf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = "✝"
r.font.name = FONT_TITLE
r.font.size = Pt(22)
r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
r.font.bold = True

# ---------- Titre principal en serif italique ----------
title_tb = slide.shapes.add_textbox(
    Inches(0.4), Inches(0.25), Inches(10), Inches(1.1)
)
ttf = title_tb.text_frame
ttf.word_wrap = True
ttf.margin_left = Inches(0.1)
p = ttf.paragraphs[0]
p.alignment = PP_ALIGN.LEFT
r = p.add_run()
r.text = "Métriques d’évaluation"
r.font.name = FONT_TITLE
r.font.size = Pt(54)
r.font.italic = True
r.font.bold = False
r.font.color.rgb = COLOR_DARK

# ---------- 4 colonnes bannière ----------
LEFT_MARGIN = Inches(0.55)
TOP_BANNERS = Inches(1.65)
N_COLS = 4
GAP = Inches(0.35)
USABLE_W = prs.slide_width - 2 * LEFT_MARGIN
col_w = (USABLE_W - GAP * (N_COLS - 1)) / N_COLS

ICON_H = Inches(0.65)
LABEL_H = Inches(0.55)
BODY_H = Inches(3.5)
RIBBON_TAIL_H = Inches(0.4)

columns = [
    {
        "icon": "📊",
        "label": "ROUGE",
        "bullets": [
            "Recouvrement",
            "n-gramme",
            "ROUGE-1",
            "ROUGE-2",
            "ROUGE-L",
        ],
    },
    {
        "icon": "🧠",
        "label": "Similarité",
        "bullets": [
            "Cosine sim.",
            "all-MiniLM-L6-v2",
            "Embeddings 384 dim.",
            "Mesure sémantique",
        ],
    },
    {
        "icon": "⏱️",
        "label": "Performance",
        "bullets": [
            "Temps de réponse",
            "Longueur réponse",
            "Taux de refus",
            "(« cannot find »)",
        ],
    },
    {
        "icon": "📈",
        "label": "Stabilité",
        "bullets": [
            "n runs / question",
            "Moyenne",
            "Écart-type",
            "Reproductibilité",
        ],
    },
]


def add_textbox(left, top, width, height, text, *, size, bold=False,
                italic=False, color=COLOR_DARK, font=FONT_BODY,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
                bullets=False):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.04)
    tf.vertical_anchor = anchor

    if isinstance(text, str):
        lines = [text]
    else:
        lines = list(text)

    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(4)
        if bullets:
            prefix = "• "
        else:
            prefix = ""
        r = p.add_run()
        r.text = prefix + line
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color
    return tb


for i, col in enumerate(columns):
    cx = LEFT_MARGIN + i * (col_w + GAP)

    # --- Icône ---
    add_textbox(
        cx, TOP_BANNERS, col_w, ICON_H, col["icon"],
        size=38, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
        font="Segoe UI Emoji",
    )

    # --- Bannière label (rectangle blanc avec liseré noir) ---
    label_top = TOP_BANNERS + ICON_H + Inches(0.05)
    label_rect = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, cx, label_top, col_w, LABEL_H
    )
    label_rect.shadow.inherit = False
    label_rect.fill.solid()
    label_rect.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    label_rect.line.color.rgb = COLOR_LINE
    label_rect.line.width = Pt(1.25)

    add_textbox(
        cx, label_top, col_w, LABEL_H, col["label"],
        size=20, bold=True, anchor=MSO_ANCHOR.MIDDLE,
        align=PP_ALIGN.CENTER, font=FONT_BODY,
    )

    # --- Corps (rectangle vertical pour la liste) ---
    body_inset = Inches(0.35)
    body_left = cx + body_inset
    body_w = col_w - 2 * body_inset
    body_top = label_top + LABEL_H + Inches(0.02)
    body_rect = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, body_left, body_top, body_w, BODY_H
    )
    body_rect.shadow.inherit = False
    body_rect.fill.background()
    body_rect.line.color.rgb = COLOR_LINE
    body_rect.line.width = Pt(1.25)

    # bullets dans le corps — alignées au rectangle, centrées verticalement
    add_textbox(
        body_left, body_top, body_w, BODY_H,
        col["bullets"],
        size=15, anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.LEFT,
        bullets=True,
    )

    # --- Queue de bannière (chevron bas, plus étroit que le corps) ---
    tail_w = body_w * 0.7
    tail_left = body_left + (body_w - tail_w) / 2
    tail_top = body_top + BODY_H
    tail = slide.shapes.add_shape(
        MSO_SHAPE.DOWN_ARROW, tail_left, tail_top, tail_w, RIBBON_TAIL_H
    )
    tail.shadow.inherit = False
    tail.fill.background()
    tail.line.color.rgb = COLOR_LINE
    tail.line.width = Pt(1.25)
    tail.text_frame.text = ""

# ---------- Save ----------
import os
out = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                   "slide_metriques_evaluation.pptx")
prs.save(out)
print(f"OK -> {out}")
